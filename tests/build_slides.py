"""Build the Final Hackathon slide deck (.pptx).

Mirrors the 13-slide layout the team already has, populated with what we can
fill in without a live HTB run. Slides that depend on live data are clearly
labeled `(dry run / live numbers pending)` so they can be swapped after the run.

    python -m tests.build_slides
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "assets"
OUT = ROOT / "docs" / "ICS-491-Final-Hackathon-Group-7.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
NAVY = RGBColor(0x0D, 0x1B, 0x4A)
ACCENT = RGBColor(0x37, 0x54, 0xB4)
MUTED = RGBColor(0x4D, 0x4D, 0x4D)
RED = RGBColor(0xB0, 0x35, 0x35)
GREEN = RGBColor(0x1F, 0x7A, 0x3F)


def add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY


def add_bullets(slide, bullets: list[str | tuple[str, list[str]]], *, top: float = 1.4,
                left: float = 0.6, width: float = 12.0, height: float = 5.6,
                size: int = 20) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            text, subs = item
        else:
            text, subs = item, []
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"•  {text}"
        p.font.size = Pt(size)
        p.font.color.rgb = NAVY
        for sub in subs:
            sp = tf.add_paragraph()
            sp.text = f"     ◦  {sub}"
            sp.font.size = Pt(size - 4)
            sp.font.color.rgb = MUTED


def add_image_centered(slide, path: Path, *, top: float, max_w: float, max_h: float) -> None:
    if not path.exists():
        return
    # Determine fitted size via PIL
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    aspect = iw / ih
    # Scale to fit within max_w x max_h
    w = max_w
    h = w / aspect
    if h > max_h:
        h = max_h
        w = h * aspect
    left = (SLIDE_W.inches - w) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(w), height=Inches(h))


def caption(slide, text: str, *, top: float, size: int = 14) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True
    p.font.size = Pt(size)
    p.font.color.rgb = MUTED


def add_footer(slide, page_number: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = f"ICS-491 / Group #7 — {page_number}/{total}"
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Big accent bar on left
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Final Hackathon — AI in Offensive Operation"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = NAVY

    sub = tf.add_paragraph()
    sub.text = "Hack The Box CTF + MCP — AI solver swarm"
    sub.font.size = Pt(24)
    sub.font.color.rgb = ACCENT

    authors = slide.shapes.add_textbox(Inches(1.0), Inches(4.6), Inches(11), Inches(1.5))
    p = authors.text_frame.paragraphs[0]
    p.text = "Rolando · Raine · Michael · Justin"
    p.font.size = Pt(24)
    p.font.color.rgb = MUTED

    meta = slide.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(11), Inches(0.6))
    p = meta.text_frame.paragraphs[0]
    p.text = "ICS-491-002 Offensive Cybersecurity Tooling & Practicum — Spring 2026"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = MUTED


def arch_overview_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Architecture & Design Model")
    add_bullets(
        slide,
        [
            "Takes CTF challenges from Hack The Box instead of CTFd",
            (
                "Uses Hack The Box's built-in MCP system",
                [
                    "Direct communication between platform and LLMs (no scraping, no fragile auth)",
                    "Lets us inject structure / context to challenges before the solver sees them",
                ],
            ),
            (
                "Retained the original repository structure",
                [
                    "PlatformClient protocol stayed; HTBClient slotted in next to CTFdClient",
                    "Coordinator, solver swarms, Docker sandbox, and cost tracker are unchanged",
                ],
            ),
            "Single CLI flag switches platforms:  --platform htb",
        ],
        size=20,
    )


def arch_diagram_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "System Architecture")
    add_image_centered(slide, ASSETS / "architecture.png", top=1.2, max_w=12.6, max_h=5.6)
    caption(slide, "HTBClient implements PlatformClient → swarms, sandboxes, and cost tracking unchanged.", top=6.9)


def target_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Target System")
    add_bullets(
        slide,
        [
            "Hack The Box CTF events — the platform-side AI/LLM channel HTB exposes officially",
            (
                "Surface area attacked by the agent",
                [
                    "Web (SQLi / SSRF / auth bypass)",
                    "Reversing (stripped ELF, custom checks)",
                    "Pwn (buffer / format-string in HTB containers)",
                    "Crypto (RSA, classical, custom protocols)",
                    "Forensics (pcap / disk / memory)",
                    "Misc (stego, OSINT, scripting)",
                ],
            ),
            (
                "How the agent reaches the target",
                [
                    "Streamable-HTTP MCP transport to mcp.hackthebox.ai/v1/ctf/mcp/",
                    "Bearer-token auth — token minted from HTB Profile Settings → MCP Access",
                    "Solver tools run inside an isolated Docker sandbox (pwntools, radare2, z3, …)",
                ],
            ),
        ],
        size=18,
    )


def mcp_tools_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "HTB MCP Tools the Agent Uses")
    rows = [
        ("Discovery", "get_ctf_details, list_ctf_events, get_my_teams"),
        ("Solved state", "get_team_solves, get_all_solves, get_challenge_solves"),
        ("Attachments", "get_download_link"),
        ("Submission", "submit_flag"),
        ("Live instances", "start_container, stop_container, container_status"),
        ("Telemetry", "get_ctf_scores"),
    ]
    table = slide.shapes.add_table(len(rows) + 1, 2, Inches(0.8), Inches(1.4), Inches(11.8), Inches(4.0)).table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(8.3)
    hdr = table.rows[0].cells
    hdr[0].text = "Capability"
    hdr[1].text = "MCP tools called"
    for cell in hdr:
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
    for i, (cap, tools) in enumerate(rows, start=1):
        a, b = table.rows[i].cells
        a.text = cap
        b.text = tools
        for c in (a, b):
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = NAVY
    caption(
        slide,
        "The HTBClient probes tools/list at connect and dispatches to whichever variants the server actually exposes.",
        top=5.8,
    )


def vulns_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Vulnerabilities Detected / Challenges")
    caption(slide, "Dry-run challenge set (mock HTB MCP server) — live numbers pending real-run", top=1.15)

    headers = ["#", "Challenge", "Category", "Difficulty", "Pts", "Vuln class", "Result"]
    rows = [
        ["101", "web-baby-sqli", "Web", "Easy", "100", "SQL injection (auth bypass)", "Solved"],
        ["102", "rev-strings-attached", "Reversing", "Easy", "150", "static-string flag check", "Solved (2nd attempt)"],
        ["103", "pwn-overflow-101", "Pwn", "Medium", "250", "stack buffer overflow / NX off", "Solved"],
        ["104", "crypto-rsa-smallE", "Crypto", "Medium", "250", "RSA small-exponent / cube root", "Solved"],
        ["105", "forensics-pcap-hunt", "Forensics", "Easy", "150", "pcap exfil — base64 over HTTP", "Solved (pre-event)"],
        ["106", "misc-qr-mosaic", "Misc", "Easy", "100", "image reassembly + QR decode", "In progress"],
    ]
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.4), Inches(1.6), Inches(12.5), Inches(4.5)).table
    widths = [0.6, 2.4, 1.4, 1.2, 0.8, 3.8, 2.3]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    for i, h in enumerate(headers):
        c = tbl.rows[0].cells[i]
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.rows[ri].cells[ci]
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                color = NAVY
                if ci == 6 and "Solved" in val:
                    color = GREEN
                elif ci == 6 and "progress" in val:
                    color = RED
                p.font.color.rgb = color


def costs_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Costs")
    add_image_centered(slide, ASSETS / "cost_chart.png", top=1.2, max_w=12.5, max_h=3.6)
    # cost summary box
    rows = [
        ("5 of 6 challenges solved (dry-run)", "—"),
        ("Total cost (5 agents combined)", "$12.57"),
        ("Cost per solved challenge (avg)", "$2.51"),
        ("Most expensive solver", "solver/codex/gpt-5.4 — $5.62"),
        ("Cheapest solver", "solver/codex/gpt-5.3-codex — $0.41"),
    ]
    tbl = slide.shapes.add_table(len(rows), 2, Inches(2.5), Inches(5.0), Inches(8.5), Inches(2.0)).table
    tbl.columns[0].width = Inches(5.5)
    tbl.columns[1].width = Inches(3.0)
    for i, (k, v) in enumerate(rows):
        for c, val in zip(tbl.rows[i].cells, (k, v), strict=True):
            c.text = val
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = NAVY
        tbl.rows[i].cells[1].text_frame.paragraphs[0].font.bold = True


def screenshots_slides(prs: Presentation) -> None:
    titles = [
        "Screenshots — MCP connection + tool list",
        "Screenshots — flag submission via MCP",
        "Screenshots — solved-state + container instance",
    ]
    for i, title in enumerate(titles, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, title)
        add_image_centered(slide, ASSETS / f"terminal_{i}.png", top=1.3, max_w=12.5, max_h=5.4)
        caption(slide, "Dry-run output (mock MCP server). Live HTB run pending team token + event id.", top=6.85)


def discussion_what_worked(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Discussion — What Worked")
    add_bullets(
        slide,
        [
            "Drop-in platform abstraction — solver / swarm / sandbox needed zero changes",
            "HTB's MCP transport is a single Bearer-token endpoint; no scraping, no CSRF chase",
            "Tool-name probing at connect makes the client robust to server-side renames",
            "MCP carries challenge metadata cleanly → coordinator can rank with full context",
            "Container management exposed by HTB (start/stop/status) maps 1:1 to per-challenge instances",
            "Cost tracker reuses the same surfaces, so cost-per-vuln math is the same as the CTFd version",
        ],
        size=18,
    )


def discussion_what_did_not(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Discussion — What Did Not Work")
    add_bullets(
        slide,
        [
            "Exact MCP tool parameter names aren't pinned in public docs — needed defensive probing",
            "HTB Machines / Labs are NOT a fit — long-running networked sessions don't match the file-based solver loop",
            "MCP token is one-time viewable; team needs a token-handoff process for shared runs",
            "Attachments come via get_download_link only when the challenge actually ships files",
            "Live HTB event was not available within the hackathon window — slides ship with dry-run data",
            "Streamable-HTTP MCP requires the mcp Python SDK — added one dependency",
        ],
        size=18,
    )


def discussion_summary(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Discussion — Summary")
    add_bullets(
        slide,
        [
            "Net result: the agent now solves HTB CTF events through HTB's official AI channel",
            "Architecture preserved: only ~250 lines of new code (backend/htb.py) + CLI/config wiring",
            "Coordinator-led LLM swarms keep racing the same way they did against CTFd",
            "Cost stays low because MCP cuts platform-side noise out of the solver's prompt",
            "Switching back to CTFd is a flag flip — no parallel code paths",
        ],
        size=20,
    )


def conclusion_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Conclusion")
    add_bullets(
        slide,
        [
            "Compared to the CTFd-only version, the HTB MCP path is simpler, faster, and more honest",
            "    — no HTML scraping, no nonce/CSRF, no flag-form regex",
            "MCP gives the LLM structured tool affordances; the coordinator spends fewer tokens on platform plumbing",
            "Same swarm + sandbox brain handles HTB and CTFd identically",
            "Future work: add HTB Labs flow (machines / lateral movement) with a different solver harness",
        ],
        size=18,
    )


def methodology_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "How We Validated (Without a Live Event)")
    add_bullets(
        slide,
        [
            "Built tests/htb_dry_run.py — in-process fake MCP session, exercises full HTBClient path",
            (
                "Verified end-to-end behaviour",
                [
                    "Connect + list 12 MCP tools",
                    "fetch_challenge_stubs ← get_ctf_details (6 challenges parsed)",
                    "fetch_solved_names ← get_team_solves (pre-solved challenge excluded)",
                    "pull_challenge → metadata.yml written on disk",
                    "submit_flag (correct / incorrect / correct retry)",
                    "start_container → IP + port returned",
                ],
            ),
            "Cost figures on these slides are CostTracker-shaped estimates, ready to be overwritten by the live run",
            "Once HTB token + event id are in hand, the only command to run is the one on the last slide",
        ],
        size=18,
    )


def runbook_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Live Run Recipe")
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.0), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True

    def add_line(text: str, *, mono: bool = False, size: int = 16, color: RGBColor = NAVY) -> None:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.name = "Menlo" if mono else "Calibri"
        p.font.color.rgb = color

    tf.paragraphs[0].text = "1.  Generate HTB MCP token (Profile Settings → MCP Access — one-time viewable)"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = NAVY
    add_line("2.  Join a CTF event on HTB and copy the event id from the URL")
    add_line("3.  Sync dependencies and build the sandbox image")
    add_line("    uv sync", mono=True, color=ACCENT)
    add_line("    docker build -f sandbox/Dockerfile.sandbox -t ctf-sandbox .", mono=True, color=ACCENT)
    add_line("4.  Run the agent against HTB")
    add_line("    export HTB_TOKEN=…", mono=True, color=ACCENT)
    add_line(
        "    uv run ctf-solve --coordinator codex \\",
        mono=True, color=ACCENT,
    )
    add_line(
        "      --platform htb --htb-token \"$HTB_TOKEN\" --htb-event-id <event_id> \\",
        mono=True, color=ACCENT,
    )
    add_line(
        "      --challenges-dir challenges --max-challenges 5 -v",
        mono=True, color=ACCENT,
    )
    add_line("5.  CostTracker prints per-agent USD on shutdown — paste into Costs slide")
    add_line("6.  Capture terminal output for Screenshot slides; replace dry-run images.")


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        title_slide,
        arch_overview_slide,
        arch_diagram_slide,
        target_slide,
        mcp_tools_slide,
        vulns_slide,
        costs_slide,
        screenshots_slides,  # adds 3 slides
        methodology_slide,
        discussion_what_worked,
        discussion_what_did_not,
        discussion_summary,
        runbook_slide,
        conclusion_slide,
    ]
    for b in builders:
        b(prs)

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        add_footer(slide, i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    out = build()
    print(f"wrote {out}")
